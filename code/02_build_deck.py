"""
Build the PowerPoint deck and the speaker-notes document.

    python code/02_build_deck.py
    # -> presentation/psi_ml_llm.pptx
    # -> presentation/SPEAKER_NOTES.md

The .pptx is a real, editable PowerPoint file: open it and restyle freely.
Speaker notes land in the notes pane of each slide as well as in the markdown
file, so you can present from either.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "presentation" / "figures"
OUT = ROOT / "presentation"

W, H = Inches(13.333), Inches(7.5)          # 16:9
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
FONT = "Helvetica Neue"

# (title, figure | None, body lines, speaker notes, kind)
#   kind: "title" | "section" | "figure" | "bullets" | "do"
DECK = [
    ("Machine Learning & Language Models", None,
     ["Teaching computers to find patterns in data"],
     "Introduce yourself briefly. Set up the hour: we will meet the four main "
     "kinds of machine learning by using each one on real cancer data, and then "
     "we will look at how a large language model works. These are two different "
     "kinds of tools, and we will treat them separately.",
     "title"),

    ("What is machine learning?", None,
     ["Finding patterns in data, without being handed the rule.",
      "",
      "You never wrote down that a cat has pointy ears and whiskers.",
      "You just saw a few thousand cats.",
      "",
      "The computer learns the same way, from examples."],
     "Keep this short. The single idea is that nobody writes the rule down. Ask "
     "the students how they would write instructions for a computer to "
     "recognize a friend's face. You cannot really do it with rules, but you "
     "can show the computer thousands of examples and let it learn the "
     "pattern, and that is what machine learning is.",
     "bullets"),

    ("Four kinds of machine learning", "ml_grid",
     [],
     "This is the map for the whole first half, so walk through it slowly. The "
     "top row is supervised learning, where we have the correct answers, called "
     "labels, and we teach the computer to reproduce them. Supervised problems "
     "split into predicting a category, which is classification, or predicting "
     "a number, which is regression. The bottom row is unsupervised learning, "
     "where there are no labels and the computer looks for structure on its "
     "own, either by grouping things (clustering) or by simplifying them "
     "(dimensionality reduction). We will do one example of each, all on the "
     "same tumor data. Come back to this grid between examples so the students "
     "keep their bearings.",
     "figure"),

    ("Where the data comes from", "mask_vs_phenotype",
     [],
     "The data comes from 41 women with triple-negative breast cancer, which is "
     "an aggressive form with no targeted treatment. A machine measured 36 "
     "proteins in every cell, which comes to about 200,000 cells. Every cell "
     "was outlined and then labeled with its type, and that labeled table is "
     "what feeds all four kinds of machine learning. We will see how the images "
     "become a table in the second example.",
     "figure"),

    ("Notebook 1: open it now", None,
     ["The four kinds of machine learning,",
      "on 20,000 real tumor cells.",
      "",
      "Every cell runs on its own. Look for the 🎛️ Try it boxes."],
     "Get everyone into notebook 1 and wait for the slowest laptop. Run the "
     "setup cell and load the cells together. Remind the students that nothing "
     "is broken or blank: they run each cell in order and then experiment with "
     "the knobs in the Try it boxes. You will drive the concepts from the "
     "slides while they follow along and play with the code.",
     "do"),

    ("① Classification: predict a category", "classification",
     [],
     "This is supervised learning. We have labels, because a biologist named "
     "every cell, so we teach the computer to reproduce them and then test it "
     "on cells it never saw. It gets about 90 percent right using two markers, "
     "and about 95 percent using all 16. The picture shows the rule it learned, "
     "which we call a decision boundary. If a cell lands in the blue region, "
     "the computer calls it a tumor cell, and that shaded boundary is the whole "
     "idea of a classifier.",
     "figure"),

    ("② Regression: predict a number", "regression",
     [],
     "This is still supervised, but the answer is a number now instead of a "
     "category. A pathologist looks at each tumor and gives it a "
     "tumor-infiltrating lymphocyte score, or TILs score, from 1 to 4, based on "
     "how much immune presence they see. The question is whether the computer's "
     "automatic immune-cell count can predict that human score. The fitted line "
     "says it can, about two thirds of the way there, with an R-squared of "
     "0.66. Regression means fitting a line and reading off a number, and this "
     "example is also the bridge from images to the clinic.",
     "figure"),

    ("③ Clustering: find groups with no labels", "cluster_heatmap",
     [],
     "Now we hide the labels, which makes this unsupervised. k-means sorted "
     "20,000 cells into 6 clusters knowing nothing about cell types. Let the "
     "students read the heatmap and name each cluster before the reveal: MPO "
     "marks neutrophils, CD20 marks B cells, CD3, CD4 and CD8 mark T cells, and "
     "the keratins mark tumor cells. The point to make is that classification "
     "reproduces labels you already have, while clustering discovers groups you "
     "did not know were there. Same data, opposite starting point.",
     "figure"),

    ("④ Dimensionality reduction: draw a map", "pca",
     [],
     "This is also unsupervised. Each cell is 16 numbers, which is a point in "
     "16-dimensional space that we cannot picture. PCA squashes those 16 "
     "numbers down to 2 while keeping the most important information, so every "
     "cell becomes a dot on a map. Tumor and immune cells land in different "
     "regions, and we add the colors afterward, so they were never used to "
     "build the map. The structure was already in the numbers. UMAP is a "
     "fancier version of the same idea that the students will run into in "
     "papers.",
     "figure"),

    ("⑤ Capstone: does composition predict survival?", "km_composition_null",
     [],
     "Now we use the toolkit on the real research question. We start with a "
     "first feature, which is the mix of cell types each patient has. We "
     "cluster the patients on their cell composition and then look at survival. "
     "The two curves sit almost on top of each other and are not statistically "
     "different, so composition on its own tells us very little about survival. "
     "That is not a dead end, because it tells us we need a better feature, "
     "which sets up feature engineering on the next slides.",
     "figure"),

    ("Same amount of immune cells, different layout", "two_patients",
     [],
     "Both of these patients have about half immune cells, which is printed on "
     "each panel, so the difference is where those immune cells sit. On the "
     "left they are walled off in their own territory, and on the right they "
     "are mixed all through the tumor. Cell composition cannot see this "
     "difference, because the counts are the same. Ask the students how they "
     "might turn walled-off versus mixed into a single number, because that "
     "number is the feature we are about to engineer.",
     "figure"),

    ("Feature engineering: where the cells are", "km_mixing",
     [],
     "This is the payoff. We engineered a new feature called the mixing score, "
     "which measures how much the immune and tumor cells touch each other, and "
     "we fed it the same survival data. Patients whose immune cells were mixed "
     "into the tumor did much worse, with about 5 times the risk of dying, and "
     "the split matches the paper's published labels for all 33 scoreable "
     "patients. The lesson is that the four kinds of machine learning are the "
     "tools, and the skill is engineering the right feature to give them. We "
     "could not have found this second feature without the first one that "
     "failed.",
     "figure"),

    ("Example 2: Bio-imaging", None,
     ["Every number we just used started as a photograph."],
     "This is the transition to the second example. All four vignettes assumed "
     "a tidy table of numbers, but somebody had to build that table from raw "
     "microscope images, and that is the job of image analysis.",
     "section"),

    ("Notebook 2: run it yourself", None,
     ["Watershed: threshold, find the centers, flood outwards.",
      "",
      "Three lines of math, and no AI at all."],
     "The students can run this one as well. The watershed method has a nice "
     "analogy: imagine each cell is a bathtub with its own shape, and we fill "
     "all the tubs with water at the same time, stopping just before they "
     "overflow, so the rising water traces the edges. In code it is three "
     "steps: threshold the bright pixels as cells, find the center of each "
     "blob, and flood outward from those centers until the floods meet. There "
     "is no neural network here, just classical image processing.",
     "do"),

    ("On easy cells, it works", "watershed_easy",
     [],
     "On tidy, well-separated cells the classical method works well. There are "
     "158 cells actually present, it finds about 177, and roughly 72 percent of "
     "the outlines are right. Old-fashioned math with no AI does the job here. "
     "Then we try it on real tissue.",
     "figure"),

    ("On real tissue, watershed struggles", "watershed_hard",
     [],
     "This is the same code on a real, densely packed tumor. Look at the big "
     "merged blobs, where the method glued neighboring cells together. It gets "
     "roughly the right count, but the shapes are wrong, and a wrong shape "
     "means protein gets attributed to the wrong cell. That error trickles down "
     "through the whole analysis, and we would not be able to recover the "
     "spatial feature from the first example. This is what motivates the next "
     "slide.",
     "figure"),

    ("So they used a neural network", None,
     ["The scientists trained a neural network on thousands of",
      "hand-drawn cell outlines, until it learned what a cell",
      "boundary looks like.",
      "",
      "It drew every one of the 200,000 outlines in our data.",
      "",
      "The same idea shows up in self-driving cars, radiation-therapy",
      "outlines, and satellite maps of forests."],
     "The takeaway is that we did not start with deep learning, we arrived at "
     "it because the simple method broke. Real tissue needed something that "
     "could learn what a boundary looks like from many examples, and that is a "
     "neural network. Segmentation like this is really just classification run "
     "once for every pixel. Give the students a familiar example or two, such "
     "as self-driving cars or outlining tumors for radiation therapy, and then "
     "move on without going deep into how neural networks work.",
     "bullets"),

    ("Example 3: Deep learning and language models", None,
     ["What is actually inside ChatGPT? Let's take one apart."],
     "This is a hard reset to the third example. Make it clear that a language "
     "model is a different kind of tool from the machine learning we just did, "
     "and not simply a bigger version of it. A large language model is a neural "
     "network built from a huge number of connected units called neurons, "
     "sometimes billions of them.",
     "section"),

    ("It starts with one neuron", "neuron",
     [],
     "A single neuron does three things: it multiplies each input by a weight, "
     "adds the results together, and squashes the total into a yes-or-no "
     "answer. That is the whole unit. In notebook 3 the students train one "
     "neuron in about ten lines to spot immune cells, which is really a tiny "
     "classifier doing the same job as the first example. A language model is "
     "this same idea repeated an enormous number of times, with billions of "
     "these weights, which we call parameters.",
     "figure"),

    ("A model cannot read letters", "tokens",
     [],
     "A model cannot read letters directly. Text is first chopped into pieces "
     "called tokens, and each token becomes a number. Have the students run "
     "their own name through the tokenizer in the notebook, which always gets a "
     "reaction. Point out that a number like 1847362 gets split into 18, 47 and "
     "362, which is a real reason these models struggle with arithmetic and "
     "with counting the letters in a word.",
     "figure"),

    ("It guesses the next token in a sequence", "next_token",
     [],
     "This is the key idea for language models. The real GPT-2 model is running "
     "on their laptop, and here it is about 95 percent sure the next token "
     "after 'diagnosed with breast' is 'cancer'. To write a sentence, it picks "
     "a token, adds it to the end, and asks the same question again. Notice "
     "that guessing the next token out of about 50,000 options is really just "
     "classification with 50,000 categories, which ties back to the first "
     "example.",
     "figure"),

    ("Temperature: it re-weights, it does not delete", "temperature_mechanism",
     [],
     "This is the knob students always ask about. To add some randomness, the "
     "model divides its scores by a number called the temperature before "
     "turning them into probabilities. A temperature below 1 stretches the "
     "gaps, so the favorite token wins more often, and 'man' goes from 45 "
     "percent up to 77 percent. A temperature above 1 shrinks the gaps, so "
     "long-shot tokens get a real chance, and 'man' drops to 28 percent. The "
     "same seven tokens are on screen the whole time, so nothing is ever "
     "removed. The knob that actually removes tokens is a different one called "
     "top-k or top-p.",
     "figure"),

    ("Turning up the temperature", "temperature_effect",
     [],
     "This is the same dial on the newer 2024 chat model, which is coherent "
     "enough that the change is easy to see. At low temperature it is fluent, "
     "at medium it is still fluent, and at high temperature it collapses into "
     "word salad. Then make the real point: read the first two answers out "
     "loud, because both sound confident and both are wrong. The sky is blue "
     "because of Rayleigh scattering, and one answer even says the sky looks "
     "white or gray. Nothing in the model was checking whether the answer was "
     "true, because a confident sentence is just a likely sentence. This is the "
     "most useful idea the students will take away, so give it time.",
     "figure"),

    ("How it keeps track of meaning", "attention",
     [],
     "This slide is optional and the first thing to cut for time. In the "
     "sentence 'The nurse examined the patient because she was worried', who is "
     "she? Inside the model, the token 'she' looks back at 'nurse'. This "
     "looking-back is called attention, and it is the T in GPT. This is one "
     "attention head out of 144, chosen because it shows the pattern clearly.",
     "figure"),

    ("The ladder", None,
     ["the neuron you built            3 numbers",
      "GPT-2  (2019)                   124,000,000",
      "the chatbot you used            500,000,000",
      "ChatGPT / Claude                ~1,000,000,000,000",
      "",
      "The same three operations all the way up: multiply, add, and squash."],
     "This closes the language-model section. The only thing that needs to land "
     "is the jump from 3 numbers to about a trillion. Nothing new appears at "
     "the top of the ladder, because it is the same arithmetic repeated an "
     "enormous number of times.",
     "bullets"),

    ("What you did today", None,
     ["You met all four kinds of machine learning, and used each one.",
      "You engineered a feature that predicts survival.",
      "You broke a classic algorithm, and saw why deep learning exists.",
      "You ran a real language model and watched it guess.",
      "",
      "All in Python, all free, and all yours to keep."],
     "End on what the students did, not on the tools. The notebooks stay "
     "theirs, so point them to the Try it boxes and the extra ideas at the "
     "bottom of each one, including asking the chatbot how many R's are in the "
     "word strawberry. Then take questions.",
     "bullets"),
]


def add_text(slide, text, left, top, width, height, size, color=INK,
             bold=False, align=PP_ALIGN.LEFT, spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = FONT
    return box


def fit_picture(slide, path, top, max_h):
    """Insert an image centerd in the box below the title.

    Wide figures end up much shorter than max_h, so centring vertically as well
    as horizontally keeps them from hugging the title with a dead band beneath.
    """
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    max_w = W - Inches(1.0)
    h, w = max_h, Emu(int(max_h * iw / ih))
    if w > max_w:
        w = max_w
        h = Emu(int(max_w * ih / iw))
    slide.shapes.add_picture(str(path), Emu(int((W - w) / 2)),
                             Emu(int(top + (max_h - h) / 2)),
                             width=w, height=h)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    for title, fig, body, notes, kind in DECK:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = SURFACE

        if kind == "title":
            add_text(slide, title, Inches(1), Inches(2.4), W - Inches(2),
                     Inches(1.6), 46, INK, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, "\n".join(body), Inches(1), Inches(4.1),
                     W - Inches(2), Inches(1.6), 22, INK2,
                     align=PP_ALIGN.CENTER, spacing=1.4)
        elif kind == "section":
            add_text(slide, title, Inches(1), Inches(2.7), W - Inches(2),
                     Inches(1.6), 38, BLUE, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, "\n".join(body), Inches(1), Inches(4.2),
                     W - Inches(2), Inches(1.2), 22, INK2,
                     align=PP_ALIGN.CENTER)
        elif kind == "do":
            add_text(slide, title, Inches(1), Inches(2.2), W - Inches(2),
                     Inches(1.2), 40, BLUE, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, "\n".join(body), Inches(1), Inches(3.7),
                     W - Inches(2), Inches(2.2), 24, INK,
                     align=PP_ALIGN.CENTER, spacing=1.5)
        elif kind == "figure":
            add_text(slide, title, Inches(0.6), Inches(0.32), W - Inches(1.2),
                     Inches(0.9), 30, INK, bold=True, align=PP_ALIGN.CENTER)
            fit_picture(slide, FIGS / f"{fig}.png", Inches(1.32),
                        H - Inches(1.85))
        else:   # bullets
            add_text(slide, title, Inches(0.9), Inches(0.6), W - Inches(1.8),
                     Inches(1.0), 34, INK, bold=True)
            add_text(slide, "\n".join(body), Inches(1.1), Inches(2.0),
                     W - Inches(2.2), Inches(4.6), 23, INK2, spacing=1.45)

        slide.notes_slide.notes_text_frame.text = notes

    path = OUT / "psi_ml_llm.pptx"
    prs.save(path)
    print(f"  {path.name}  ({len(DECK)} slides)")
    return len(DECK)


TIMING = """# Speaker notes: Machine Learning & Language Models

**Audience:** high school, no coding assumed.
**Nominal length:** about 60 minutes. The deck is built a little long on
purpose, so the cut list below tells you what to drop if you run behind.

## Shape of the hour

| min | what |
|-----|------|
| 0 to 6 | What machine learning is, plus the four-kinds grid (slides 1 to 4) |
| 6 to 10 | Open notebook 1 (slide 5) |
| 10 to 30 | **Notebook 1**, the four examples: classify, regress, cluster, PCA (slides 6 to 9) |
| 30 to 38 | **Capstone**, engineer a feature and predict survival (slides 10 to 12) |
| 38 to 46 | **Notebook 2**, image segmentation (slides 13 to 17) |
| 46 to 58 | **Notebook 3**, language models (slides 18 to 24) |
| 58 to 60 | Wrap (slides 25 to 26) |

## Cut list, in this order, if you are running late

1. **Slide 9 (PCA / dimensionality reduction).** It has the smallest payoff of
   the four examples, so you can name it and move on. Saves about 3 minutes.
2. **Slide 24 (attention).** The language-model story holds together without it.
   Saves about 4 minutes.
3. **Notebook 3, the numpy neuron.** Jump from the neuron diagram straight to
   tokenization. Saves about 4 minutes.
4. **Notebook 2 as slides only.** Show watershed_easy and watershed_hard instead
   of running the code live. Saves about 5 minutes.
5. **Last resort, the capstone (slides 10 to 12).** This is the best science in
   the hour and the biggest payoff, so cut it only if you truly must, and if you
   do, at least show slide 12 (km_mixing) as a still image.

## Before class

- [ ] The notebooks currently point at a **local** data path for testing
      (`DATA = "../data/processed"`). Before you share them with students,
      switch each notebook's `DATA` line back to the GitHub URL, which is kept
      right above it as a comment. Otherwise Colab will not find the data.
- [ ] Open notebook 3 once on the day, because the first GPT-2 download is the
      slowest step and Colab does not cache it between sessions.
- [ ] With the GitHub URL active, run `python code/04_verify.py --live` and
      confirm every check passes.

## If the wifi dies

Notebooks 1 and 2 need only a 4 MB download, while notebook 3 also needs about
1.5 GB of model weights. If the network is bad, present the language-model
section from the slides, since every figure there is a real output from the
models, so you lose only the live interaction.

## Numbers you will be asked about

- **Classification:** about 89% on 2 markers, and **95%** on all 16 (random
  forest, 3-fold cross-validation).
- **Regression:** immune fraction versus pathologist TILs score, **R^2 = 0.66**,
  n = 25.
- 41 patients were imaged, and **38** have both cell data and clinical follow-up.
- **33** patients are scoreable; the other 5 are "cold" (fewer than 250 immune
  cells) and are set aside.
- Mixed versus walled off: **hazard ratio 5.21, p = 0.032** (log-rank p = 0.017).
- The original paper reported HR 4.97, p = 0.03. We are reproducing the finding,
  not matching it to the decimal, because our cell-contact rule is simpler.
- Survival here is **overall survival**, not disease-free survival, so say
  "survival".

## Honest caveats worth saying out loud

- 38 patients is a *small* study. One or two patients moving could change the
  p-value, and a real conclusion would need hundreds.
- We chose the 0.26 cutoff after looking at the data. Tell the students that,
  and then point them at the notebook box that tries other cutoffs.
- The segmentation images in notebook 2 are **simulated** from the real cell
  outlines, because the raw microscope channels are not in the public download.
  The difficulty is real, but the pixels are not.

---

## Per-slide script

"""


def build_notes():
    lines = [TIMING]
    for i, (title, fig, body, notes, kind) in enumerate(DECK, 1):
        lines.append(f"### {i}. {title}")
        if fig:
            lines.append(f"*figure: `{fig}.png`*\n")
        if body:
            lines.append("> " + "\n> ".join(b for b in body if b) + "\n")
        lines.append(notes + "\n")
    path = OUT / "SPEAKER_NOTES.md"
    path.write_text("\n".join(lines))
    print(f"  {path.name}")


if __name__ == "__main__":
    print("presentation:")
    build()
    build_notes()
